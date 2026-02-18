#!/usr/bin/env python3
"""Generate Offline vs Online Comparison PowerPoint with embedded images."""

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BG_COLOR = RGBColor(0x1B, 0x1B, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x00, 0xBC, 0xD4)
TEAL = RGBColor(0x00, 0x96, 0x88)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_CELL = RGBColor(0x2A, 0x2A, 0x45)
HEADER_BG = RGBColor(0x00, 0x77, 0x88)
MARGIN = Inches(0.3)

IMG_DIR = str(PROJECT_ROOT / "results")
IMG_BASE_TRAIN = f"{IMG_DIR}/base_model_offline/base_model_training_progress.png"
IMG_BASE_FLOW = f"{IMG_DIR}/base_model_offline/base_flow_visualization.png"
IMG_LARGE_TRAIN = f"{IMG_DIR}/large_model_offline/large_model_training_progress.png"
IMG_LARGE_FLOW = f"{IMG_DIR}/large_model_offline/large_flow_visualization.png"
IMG_BASE_ONLINE = f"{IMG_DIR}/base_model_online/base_online_visualization.png"
IMG_LARGE_ONLINE_PROG = f"{IMG_DIR}/large_model_online/large_model_online_progress.png"
IMG_LARGE_ONLINE_VIS = f"{IMG_DIR}/large_model_online/large_online_visualization.png"

OUT_PATH = str(PROJECT_ROOT / "Offline_vs_Online_Comparison.pptx")


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_title(slide, text, font_size=32):
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
                text, font_size=font_size, bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT)


def set_cell(cell, text, font_size=14, bold=False, color=WHITE, bg=DARK_CELL, alignment=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cf = cell.fill
    cf.solid()
    cf.fore_color.rgb = bg


def add_full_image(slide, img_path, top_offset=Inches(1.1)):
    from PIL import Image
    avail_w = SLIDE_W - 2 * MARGIN
    avail_h = SLIDE_H - top_offset - MARGIN
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    if avail_w / avail_h > aspect:
        h = avail_h
        w = int(h * aspect)
    else:
        w = avail_w
        h = int(w / aspect)
    left = (SLIDE_W - w) // 2
    slide.shapes.add_picture(img_path, left, top_offset, w, h)


def add_side_by_side_images(slide, img_left_path, img_right_path,
                             label_left="", label_right="",
                             top_offset=Inches(1.1)):
    img_w = Inches(6.0)
    img_h = Inches(4.5)
    gap = Inches(0.333)
    total_w = 2 * img_w + gap
    start_left = (SLIDE_W - total_w) // 2

    slide.shapes.add_picture(img_left_path, start_left, top_offset, img_w, img_h)
    slide.shapes.add_picture(img_right_path, start_left + img_w + gap, top_offset, img_w, img_h)

    label_top = top_offset + img_h + Inches(0.05)
    add_textbox(slide, start_left, label_top, img_w, Inches(0.35),
                label_left, font_size=14, bold=False, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_left + img_w + gap, label_top, img_w, Inches(0.35),
                label_right, font_size=14, bold=False, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)


# Build Presentation
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# SLIDE 1 - Title
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_textbox(sl, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
            "Offline vs Online Training Comparison",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
            "Neural Network-Based Flow Field Compression",
            font_size=24, bold=False, color=CYAN, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
            "Mahesh Sadupalli",
            font_size=20, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = CYAN
shape.line.fill.background()

# SLIDE 2 - Compression Ratio Analysis
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Compression Ratio Analysis")

txBox = sl.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
lines = [
    ("Original Data: 7,919,100 rows \u00d7 8 columns (4 input + 4 output), float32", WHITE, 18, True),
    ("Model Sizes:  Base .pth = 30 KB  |  Large .pth = 104 KB", LIGHT_GRAY, 16, False),
]
for i, (txt, clr, sz, bld) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(sz)
    p.font.bold = bld
    p.font.color.rgb = clr
    p.font.name = "Calibri"

rows, cols = 3, 4
tbl = sl.shapes.add_table(rows, cols, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.8)).table
tbl.columns[0].width = Inches(3.5)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.4)
tbl.columns[3].width = Inches(2.4)
headers = ["Method", "Size", "Base CR", "Large CR"]
data = [
    ["Method 1: Binary (.bin)", "253 MB", "8,277 : 1", "2,379 : 1"],
    ["Method 2: CSV raw (.csv)", "833 MB", "27,208 : 1", "7,819 : 1"],
]
for c, h in enumerate(headers):
    set_cell(tbl.cell(0, c), h, font_size=15, bold=True, color=WHITE, bg=HEADER_BG)
for r, row in enumerate(data, 1):
    for c, val in enumerate(row):
        clr = GREEN if ":" in val else WHITE
        set_cell(tbl.cell(r, c), val, font_size=14, color=clr, bg=DARK_CELL)

# SLIDE 3 - Offline Training Results
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Offline Training Results")

rows, cols = 3, 6
tbl = sl.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.0)).table
tbl.columns[0].width = Inches(1.5)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(1.7)
tbl.columns[3].width = Inches(1.8)
tbl.columns[4].width = Inches(1.8)
tbl.columns[5].width = Inches(1.9)
headers = ["Model", "Architecture", "Parameters", "PSNR (dB)", "SSIM", "Rel. Error (%)"]
data = [
    ["Base",  "4\u219264\u219264\u219232\u21924",    "6,692",  "30.90", "0.920", "5.29"],
    ["Large", "4\u2192128\u2192128\u219264\u21924", "25,668", "35.59", "0.980", "3.04"],
]
for c, h in enumerate(headers):
    set_cell(tbl.cell(0, c), h, font_size=14, bold=True, color=WHITE, bg=HEADER_BG)
for r, row in enumerate(data, 1):
    for c, val in enumerate(row):
        clr = GREEN if c in (3, 4) else (RED if c == 5 else WHITE)
        set_cell(tbl.cell(r, c), val, font_size=14, color=clr, bg=DARK_CELL)

# SLIDE 4 - Offline Training Progress (side-by-side)
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Offline Training Progress")
add_side_by_side_images(sl, IMG_BASE_TRAIN, IMG_LARGE_TRAIN,
                         "Base Model (64-64-32)", "Large Model (128-128-64)")

# SLIDE 5 - Base Offline Flow
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Base Model \u2014 Flow Field Reconstruction (Offline)")
add_full_image(sl, IMG_BASE_FLOW)

# SLIDE 6 - Large Offline Flow
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Large Model \u2014 Flow Field Reconstruction (Offline)")
add_full_image(sl, IMG_LARGE_FLOW)

# SLIDE 7 - Online Streaming Results
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Online Streaming Results")
add_textbox(sl, Inches(0.5), Inches(1.05), Inches(12), Inches(0.5),
            "Training metrics vs Evaluation on full dataset \u2014 reveals catastrophic forgetting",
            font_size=16, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

rows, cols = 3, 5
tbl = sl.shapes.add_table(rows, cols, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.4)).table
tbl.columns[0].width = Inches(2.3)
tbl.columns[1].width = Inches(2.25)
tbl.columns[2].width = Inches(2.25)
tbl.columns[3].width = Inches(2.25)
tbl.columns[4].width = Inches(2.25)
headers = ["Phase", "Base PSNR", "Base SSIM", "Large PSNR", "Large SSIM"]
for c, h in enumerate(headers):
    set_cell(tbl.cell(0, c), h, font_size=14, bold=True, color=WHITE, bg=HEADER_BG)
train_row = ["Training (last window)", "23.85", "0.890", "24.92", "0.894"]
eval_row = ["Evaluation (full data)", "11.52", "0.758", "6.42", "0.684"]
for c, val in enumerate(train_row):
    clr = CYAN if c > 0 else WHITE
    set_cell(tbl.cell(1, c), val, font_size=14, color=clr, bg=DARK_CELL)
for c, val in enumerate(eval_row):
    clr = RED if c > 0 else WHITE
    set_cell(tbl.cell(2, c), val, font_size=14, color=clr, bg=DARK_CELL)

shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.7))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x44, 0x1C, 0x1C)
shape.line.color.rgb = RED
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "\u26a0  Catastrophic Forgetting: Larger models degrade MORE on full evaluation"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RED
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# SLIDE 8 - Online Training Progress (Large)
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Online Training Progress \u2014 Large Model")
add_full_image(sl, IMG_LARGE_ONLINE_PROG)
add_textbox(sl, Inches(1), Inches(6.9), Inches(11.3), Inches(0.4),
            "PSNR rises then collapses as model forgets earlier windows",
            font_size=13, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# SLIDE 9 - Base Online Flow
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Base Model \u2014 Flow Field Reconstruction (Online)")
add_full_image(sl, IMG_BASE_ONLINE)

# SLIDE 10 - Large Online Flow
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Large Model \u2014 Flow Field Reconstruction (Online)")
add_full_image(sl, IMG_LARGE_ONLINE_VIS)

# SLIDE 11 - Head-to-Head Comparison
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Offline vs Online \u2014 Head-to-Head")

rows, cols = 3, 5
tbl = sl.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.0)).table
tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(2.375)
tbl.columns[2].width = Inches(2.375)
tbl.columns[3].width = Inches(2.375)
tbl.columns[4].width = Inches(2.375)
headers = ["Metric", "Base Offline", "Base Online", "Large Offline", "Large Online"]
for c, h in enumerate(headers):
    set_cell(tbl.cell(0, c), h, font_size=14, bold=True, color=WHITE, bg=HEADER_BG)

psnr_vals = ["PSNR (dB)", "30.90", "11.52", "35.59", "6.42"]
ssim_vals = ["SSIM",       "0.920", "0.758", "0.980", "0.684"]

for c, val in enumerate(psnr_vals):
    if c == 0:
        clr = WHITE
    elif c in (1, 3):
        clr = GREEN
    else:
        clr = RED
    set_cell(tbl.cell(1, c), val, font_size=15, bold=(c > 0), color=clr, bg=DARK_CELL)

for c, val in enumerate(ssim_vals):
    if c == 0:
        clr = WHITE
    elif c in (1, 3):
        clr = GREEN
    else:
        clr = RED
    set_cell(tbl.cell(2, c), val, font_size=15, bold=(c > 0), color=clr, bg=DARK_CELL)

# Visual bar comparison
bar_data = [
    ("Base Offline  PSNR 30.90", 30.90, GREEN),
    ("Base Online   PSNR 11.52", 11.52, RED),
    ("Large Offline PSNR 35.59", 35.59, GREEN),
    ("Large Online  PSNR  6.42",  6.42, RED),
]
bar_top = Inches(4.5)
bar_left = Inches(1.5)
max_bar_w = Inches(8.0)
max_val = 36.0
for i, (label, val, clr) in enumerate(bar_data):
    y = bar_top + Inches(i * 0.6)
    w = int(max_bar_w * (val / max_val))
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y, w, Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    add_textbox(sl, bar_left + w + Inches(0.15), y, Inches(3.5), Inches(0.4),
                label, font_size=12, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

# SLIDE 12 - Key Takeaways
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl)
add_title(sl, "Key Takeaways")

bullets = [
    ("Offline training achieves excellent compression quality (PSNR >30 dB, SSIM >0.92)", GREEN),
    ("Compression ratios of 8,277:1 (Base) and 2,379:1 (Large) using binary method", GREEN),
    ("Online streaming suffers from catastrophic forgetting \u2014 final model only remembers recent windows", RED),
    ("Larger networks are MORE susceptible to catastrophic forgetting in online mode", RED),
    ("Online training metrics are misleading \u2014 evaluation on full data reveals the true performance gap", RED),
    ("Mitigation strategies needed: experience replay, EWC, or progressive networks", CYAN),
]

txBox = sl.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True

for i, (bullet, marker_color) in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    run1 = p.add_run()
    run1.text = "\u25b6  "
    run1.font.size = Pt(16)
    run1.font.color.rgb = marker_color
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = bullet
    run2.font.size = Pt(18)
    run2.font.color.rgb = WHITE
    run2.font.name = "Calibri"

# Save
prs.save(OUT_PATH)
print(f"Presentation saved to {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
