"""
Architecture diagrams — ONLY clean diagrams, no text outside safe margins.
Two slides: Offline INR architecture, Online temporal window architecture.
All native PPT elements, fully editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = 12192000
prs.slide_height = 6858000

# ─── Colors ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)

BLUE_F = RGBColor(0x00, 0x6E, 0xB8)
BLUE_T = RGBColor(0x4D, 0xA6, 0xE8)
BLUE_S = RGBColor(0x00, 0x4C, 0x8C)
ORANGE_F = RGBColor(0xE8, 0x8D, 0x2A)
ORANGE_T = RGBColor(0xFF, 0xB3, 0x47)
ORANGE_S = RGBColor(0xCC, 0x77, 0x22)

LAYER_GREEN = RGBColor(0x2E, 0x8B, 0x57)
LAYER_PURPLE = RGBColor(0x6A, 0x0D, 0xAD)
RELU_COLOR = RGBColor(0xFF, 0x69, 0x34)
RED = RGBColor(0xDC, 0x35, 0x45)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
LOSS_BG = RGBColor(0xFF, 0xF8, 0xE1)
LOSS_BORDER = RGBColor(0xE8, 0x8D, 0x2A)
CALLOUT_BG = RGBColor(0xFD, 0xEB, 0xD0)
WARN_BG = RGBColor(0xFD, 0xEC, 0xEA)


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, l, t, w, h, fill, text='', sz=10, tc=WHITE, bold=True, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if text:
        s.text_frame.word_wrap = True
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = s.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = tc
        p.font.name = 'Calibri'
    return s


def rbox(slide, l, t, w, h, fill, text='', sz=10, tc=WHITE, bold=True, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    if text:
        s.text_frame.word_wrap = True
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = s.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = tc
        p.font.name = 'Calibri'
    return s


def label(slide, l, t, w, h, text, sz=10, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tb


def arrow(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, Inches(0.3))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    s.line.fill.background()
    return s


def hline(slide, l, t, w, color=DARK):
    box(slide, l, t, w, Pt(1.5), color)


def vline(slide, l, t, h, color=DARK):
    box(slide, l, t, Pt(1.5), h, color)


def bars(slide, cx, cy, n, h, color=LAYER_GREEN, bw=Inches(0.1), gap=Inches(0.04)):
    """Draw n vertical bars centered at (cx, cy). Returns (left_edge, right_edge)."""
    total = n * bw + (n - 1) * gap
    xs = cx - total // 2
    ys = cy - h // 2
    for i in range(n):
        bx = xs + i * (bw + gap)
        box(slide, bx, ys, bw, h, color)
    return xs, xs + total


def bracket_top(slide, x1, x2, y, text, sz=9):
    """Bracket above with label."""
    bh = Inches(0.12)
    mid = (x1 + x2) // 2
    vline(slide, x1, y - bh, bh)
    vline(slide, x2, y - bh, bh)
    hline(slide, x1, y - bh, x2 - x1)
    vline(slide, mid, y - bh - Inches(0.06), Inches(0.06))
    label(slide, x1 - Inches(0.5), y - bh - Inches(0.35), x2 - x1 + Inches(1.0), Inches(0.25),
          text, sz=sz, bold=True)


def bracket_bot(slide, x1, x2, y, text, sz=8):
    """Bracket below with label."""
    bh = Inches(0.1)
    vline(slide, x1, y, bh)
    vline(slide, x2, y, bh)
    hline(slide, x1, y + bh, x2 - x1)
    label(slide, x1 - Inches(0.2), y + bh + Inches(0.02), x2 - x1 + Inches(0.4), Inches(0.25),
          text, sz=sz, bold=True)


def add_table(slide, left, top, w, h, data, hdr_color, hl_cols=None, fsz=8):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, w, h)
    tbl = ts.table
    cw = w // cols
    for c in range(cols):
        tbl.columns[c].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(fsz)
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = (r == 0)
                p.font.color.rgb = WHITE if r == 0 else (RED if (hl_cols and c in hl_cols) else DARK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hdr_color if r == 0 else (WHITE if r % 2 else RGBColor(0xF5, 0xF5, 0xF5))


# ═══════════════════════════════════════════════
# SLIDE 1: OFFLINE INR ARCHITECTURE
# ═══════════════════════════════════════════════
s = blank()
box(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)

cy = Inches(3.0)  # vertical center for network

# ── INPUT CUBE (fake 3D: front + top parallelogram + side) ──
cx, ctop, cw, ch = Inches(0.5), Inches(1.6), Inches(1.8), Inches(2.2)
d = Inches(0.18)

# Top face
tp = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, cx, ctop - d, cw, d + Pt(2))
tp.fill.solid()
tp.fill.fore_color.rgb = BLUE_T
tp.line.color.rgb = BLUE_S
tp.line.width = Pt(0.8)

# Side face
box(s, cx + cw, ctop, d, ch, BLUE_S, border=BLUE_S)

# Front face
box(s, cx, ctop, cw, ch, BLUE_F, 'Input\nData', 13, WHITE, border=BLUE_S)

# Input table
add_table(s, Inches(0.3), Inches(4.1), Inches(2.2), Inches(0.7),
          [['X', 'Y', 'Z', 'Time'],
           ['1.01', '2.01', '5.02', '0.01'],
           ['N...', 'N...', 'N...', 'N...']],
          BLUE_F, fsz=7)

# ── ARROW ──
arrow(s, Inches(2.65), cy - Inches(0.15), Inches(0.55))

# ── LAYER BARS ──
# Input layer (4)
l1l, l1r = bars(s, Inches(3.75), cy, 3, Inches(0.9), LAYER_PURPLE)
# Hidden 1 (64)
l2l, l2r = bars(s, Inches(4.85), cy, 6, Inches(2.4), LAYER_GREEN)
# Hidden 2 (64)
l3l, l3r = bars(s, Inches(6.15), cy, 6, Inches(2.4), LAYER_GREEN)
# Hidden 3 (32)
l4l, l4r = bars(s, Inches(7.45), cy, 5, Inches(1.7), LAYER_GREEN)
# Output layer (4)
l5l, l5r = bars(s, Inches(8.55), cy, 3, Inches(0.9), LAYER_PURPLE)

# Small arrows between groups
for x in [l1r, l2r, l3r, l4r]:
    arrow(s, x + Inches(0.03), cy - Inches(0.1), Inches(0.28))

# Layer size labels
for xc, lbl in [(Inches(3.75), '4'), (Inches(4.85), '64'), (Inches(6.15), '64'),
                 (Inches(7.45), '32'), (Inches(8.55), '4')]:
    label(s, xc - Inches(0.2), cy + Inches(1.35), Inches(0.4), Inches(0.25),
          lbl, sz=11, bold=True)

# Bracket above hidden layers
bracket_top(s, l2l, l4r, cy - Inches(1.3), 'Hidden Layers: ReLU Activation', sz=10)

# Brackets below groups
bracket_bot(s, l1l, l1r, cy + Inches(0.5), 'Input')
bracket_bot(s, l2l, l3r, cy + Inches(1.25), 'Hidden Layers (ReLU)')
bracket_bot(s, l4l, l4r, cy + Inches(0.9), 'Hidden')
bracket_bot(s, l5l, l5r, cy + Inches(0.5), 'Output')

# ReLU labels
for rx in [Inches(4.25), Inches(5.45), Inches(6.75)]:
    label(s, rx, cy - Inches(1.35), Inches(0.5), Inches(0.2),
          'ReLU', sz=7, bold=True, color=RELU_COLOR)

label(s, Inches(7.95), cy - Inches(1.3), Inches(0.8), Inches(0.18),
      'No Activation', sz=6, color=GRAY)

# ── ARROW ──
arrow(s, l5r + Inches(0.1), cy - Inches(0.15), Inches(0.55))

# ── OUTPUT CUBE ──
ox = Inches(9.7)
tp2 = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, ox, ctop - d, cw, d + Pt(2))
tp2.fill.solid()
tp2.fill.fore_color.rgb = ORANGE_T
tp2.line.color.rgb = ORANGE_S
tp2.line.width = Pt(0.8)

box(s, ox + cw, ctop, d, ch, ORANGE_S, border=ORANGE_S)
box(s, ox, ctop, cw, ch, ORANGE_F, 'Predicted\nData', 13, WHITE, border=ORANGE_S)

# Output table
add_table(s, Inches(9.5), Inches(4.1), Inches(2.2), Inches(0.7),
          [['\u0056\u0302x', '\u0056\u0302y', '\u0050\u0302', 'TKE\u0302'],
           ['1.05', '-0.32', '0.95', '12.3'],
           ['N...', 'N...', 'N...', 'N...']],
          ORANGE_F, hl_cols=[0, 1, 2, 3], fsz=7)

# ── LOSS BOX ──
rbox(s, Inches(4.0), Inches(5.2), Inches(4.0), Inches(0.55),
     LOSS_BG, 'Loss: MSE  |  Metrics: PSNR, SSIM', 9, DARK, True, LOSS_BORDER)

# ── CALLOUT ──
rbox(s, Inches(1.5), Inches(6.0), Inches(9.0), Inches(0.45),
     CALLOUT_BG,
     '4 \u2192 64 \u2192 64 \u2192 32 \u2192 4  |  6,692 params  |  26.1 KB  |  CR ~4,734:1',
     9, RGBColor(0x8B, 0x40, 0x00), True, RGBColor(0xE8, 0x8D, 0x2A))


# ═══════════════════════════════════════════════
# SLIDE 2: ONLINE ARCHITECTURE
# Same layout as offline: Input Cube → [Windows] → [Hidden] → [Output] → Output Cube
# Plus catastrophic forgetting element below
# ═══════════════════════════════════════════════
s = blank()
box(s, 0, 0, prs.slide_width, prs.slide_height, LIGHT_BG)

cy = Inches(2.6)  # slightly higher to make room for forgetting element

# ── INPUT CUBE ──
cx, ctop, cw, ch = Inches(0.4), Inches(1.3), Inches(1.7), Inches(2.0)
d = Inches(0.15)

tp = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, cx, ctop - d, cw, d + Pt(2))
tp.fill.solid()
tp.fill.fore_color.rgb = BLUE_T
tp.line.color.rgb = BLUE_S
tp.line.width = Pt(0.8)

box(s, cx + cw, ctop, d, ch, BLUE_S, border=BLUE_S)
box(s, cx, ctop, cw, ch, BLUE_F, 'Streaming\nData', 12, WHITE, border=BLUE_S)

# Input table
add_table(s, Inches(0.3), Inches(3.55), Inches(2.0), Inches(0.6),
          [['X', 'Y', 'Z', 'Time'],
           ['1.01', '2.01', '5.02', '0.01'],
           ['N...', 'N...', 'N...', 'N...']],
          BLUE_F, fsz=6)

# ── ARROW ──
arrow(s, Inches(2.3), cy - Inches(0.12), Inches(0.4))

# ── TEMPORAL WINDOWING AREA (light blue background box) ──
win_area_x = Inches(2.85)
win_area_w = Inches(2.6)
win_area_y = Inches(1.15)
win_area_h = Inches(2.8)
rbox(s, win_area_x, win_area_y, win_area_w, win_area_h,
     RGBColor(0xE3, 0xF2, 0xFD), '', border=RGBColor(0x90, 0xCA, 0xF9))

# "Temporal Windowing" header
label(s, win_area_x, win_area_y + Inches(0.05), win_area_w, Inches(0.25),
      'Temporal Windowing', 9, True, DARK)

# Window boxes (matching progress_update style: 0.42" x 0.85")
win_w = Inches(0.42)
win_h = Inches(0.85)
win_gap = Inches(0.1)
win_row_x = win_area_x + Inches(0.12)
win_row_y = win_area_y + Inches(0.4)

win_data = [
    ('W1',  RGBColor(0x90, 0xCA, 0xF9), DARK),
    ('W2',  RGBColor(0x42, 0xA5, 0xF5), WHITE),
    ('...', RGBColor(0xF5, 0xF5, 0xF5), GRAY),
    ('W19', RGBColor(0x1E, 0x88, 0xE5), WHITE),
    ('W20', RGBColor(0x15, 0x65, 0xC0), WHITE),
]
# Fit 5 windows in the area
actual_gap = (win_area_w - Inches(0.24) - 5 * win_w) // 4

for i, (wlbl, wfill, wtc) in enumerate(win_data):
    wx = win_row_x + i * (win_w + actual_gap)
    box(s, wx, win_row_y, win_w, win_h, wfill, wlbl, 7 if len(wlbl) > 2 else 8, wtc, True,
        RGBColor(0x0D, 0x47, 0xA1) if wlbl != '...' else RGBColor(0xCC, 0xCC, 0xCC))

# "15 timesteps per window" below windows
label(s, win_area_x, win_row_y + win_h + Inches(0.08), win_area_w, Inches(0.2),
      '15 timesteps per window', 7, False, GRAY)

# "100 epochs / window" below that
label(s, win_area_x, win_row_y + win_h + Inches(0.28), win_area_w, Inches(0.2),
      '20 windows \u00d7 100 epochs = 2,000 total', 7, True, RGBColor(0x15, 0x65, 0xC0))

# ── ARROW: Windowing → Hidden layers ──
arrow(s, win_area_x + win_area_w + Inches(0.05), cy - Inches(0.12), Inches(0.35))

# ── HIDDEN LAYER BARS ──
l2l, l2r = bars(s, Inches(6.5), cy, 6, Inches(2.2), LAYER_GREEN)
l3l, l3r = bars(s, Inches(7.5), cy, 6, Inches(2.2), LAYER_GREEN)
l4l, l4r = bars(s, Inches(8.5), cy, 5, Inches(1.5), LAYER_GREEN)

# Output layer bars
l5l, l5r = bars(s, Inches(9.35), cy, 3, Inches(0.8), LAYER_PURPLE)

# Small arrows between groups
for x in [l2r, l3r, l4r]:
    arrow(s, x + Inches(0.03), cy - Inches(0.1), Inches(0.2))

# Layer size labels
for xc, lbl in [(Inches(6.5), '64'), (Inches(7.5), '64'),
                 (Inches(8.5), '32'), (Inches(9.35), '4')]:
    label(s, xc - Inches(0.2), cy + Inches(1.2), Inches(0.4), Inches(0.25),
          lbl, sz=10, bold=True)

# Bracket above hidden layers
bracket_top(s, l2l, l4r, cy - Inches(1.2), 'Hidden Layers: ReLU', sz=9)

# Brackets below
bracket_bot(s, win_area_x, win_area_x + win_area_w, cy + Inches(0.6), 'Input (4)')
bracket_bot(s, l2l, l3r, cy + Inches(1.15), 'Hidden (ReLU)')
bracket_bot(s, l4l, l4r, cy + Inches(0.8), 'Hidden')
bracket_bot(s, l5l, l5r, cy + Inches(0.45), 'Output')

# ReLU labels
for rx in [Inches(6.9), Inches(7.95)]:
    label(s, rx, cy - Inches(1.2), Inches(0.45), Inches(0.18),
          'ReLU', sz=7, bold=True, color=RELU_COLOR)

label(s, Inches(8.85), cy - Inches(1.15), Inches(0.7), Inches(0.16),
      'No Activation', sz=6, color=GRAY)

# ── ARROW: Network → Output ──
arrow(s, l5r + Inches(0.08), cy - Inches(0.12), Inches(0.4))

# ── OUTPUT CUBE ──
ox = Inches(10.2)
tp2 = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, ox, ctop - d, cw, d + Pt(2))
tp2.fill.solid()
tp2.fill.fore_color.rgb = ORANGE_T
tp2.line.color.rgb = ORANGE_S
tp2.line.width = Pt(0.8)

box(s, ox + cw, ctop, d, ch, ORANGE_S, border=ORANGE_S)
box(s, ox, ctop, cw, ch, ORANGE_F, 'Predicted\nData', 12, WHITE, border=ORANGE_S)

# Output table
add_table(s, Inches(10.1), Inches(3.55), Inches(2.0), Inches(0.6),
          [['\u0056\u0302x', '\u0056\u0302y', '\u0050\u0302', 'TKE\u0302'],
           ['1.05', '-0.32', '0.95', '12.3'],
           ['N...', 'N...', 'N...', 'N...']],
          ORANGE_F, hl_cols=[0, 1, 2, 3], fsz=6)

# ── LOSS BOX ──
rbox(s, Inches(4.8), Inches(4.35), Inches(3.5), Inches(0.45),
     LOSS_BG, 'Loss: MSE  |  Metrics: PSNR, SSIM', 8, DARK, True, LOSS_BORDER)

# ══════════════════════════════════════
# CATASTROPHIC FORGETTING ELEMENT
# ══════════════════════════════════════
forget_y = Inches(5.1)

# Red box background
rbox(s, Inches(0.4), forget_y, Inches(12.0), Inches(1.8),
     RGBColor(0xFD, 0xEC, 0xEA), '', border=RED)

# Title
label(s, Inches(0.6), forget_y + Inches(0.05), Inches(11.5), Inches(0.3),
      '\u26A0  Catastrophic Forgetting', 12, True, RED)

# Visual: W1 through W20 showing which are forgotten/remembered
fg_win_y = forget_y + Inches(0.45)
fg_win_w = Inches(0.42)
fg_win_h = Inches(0.5)
n_wins = 10  # show W1..W5, dots, W18..W20
fg_gap = Inches(0.08)
fg_start_x = Inches(0.7)

fg_windows = [
    ('W1', RGBColor(0xEF, 0x9A, 0x9A), RED, True),     # faded red = forgotten
    ('W2', RGBColor(0xEF, 0x9A, 0x9A), RED, True),
    ('W3', RGBColor(0xEF, 0x9A, 0x9A), RED, True),
    ('W4', RGBColor(0xEF, 0x9A, 0x9A), RED, True),
    ('W5', RGBColor(0xEF, 0x9A, 0x9A), RED, True),
    ('...', RGBColor(0xF5, 0xF5, 0xF5), GRAY, False),
    ('W18', RGBColor(0xEF, 0x9A, 0x9A), RED, True),
    ('W19', RGBColor(0xEF, 0x9A, 0x9A), RED, True),
    ('W20', RGBColor(0xA5, 0xD6, 0xA7), GREEN, False),  # green = remembered
]

for i, (wlbl, wfill, wtc, forgotten) in enumerate(fg_windows):
    wx = fg_start_x + i * (fg_win_w + fg_gap)
    box(s, wx, fg_win_y, fg_win_w, fg_win_h, wfill, wlbl, 7, wtc, True,
        RED if forgotten else GREEN)

    # Red X for forgotten, green check for remembered
    if wlbl != '...':
        sym = '\u2717' if forgotten else '\u2713'
        sym_color = RED if forgotten else GREEN
        label(s, wx, fg_win_y + fg_win_h + Pt(2), fg_win_w, Inches(0.18),
              sym, 10, True, sym_color)

# Text explanation on right
explain_x = fg_start_x + 9 * (fg_win_w + fg_gap) + Inches(0.2)
label(s, explain_x, fg_win_y - Inches(0.05), Inches(4.0), Inches(0.25),
      'After training on W20:', 9, True, DARK, PP_ALIGN.LEFT)
label(s, explain_x, fg_win_y + Inches(0.2), Inches(4.0), Inches(0.2),
      '\u2022 W1\u2013W19: knowledge overwritten', 8, False, RED, PP_ALIGN.LEFT)
label(s, explain_x, fg_win_y + Inches(0.4), Inches(4.0), Inches(0.2),
      '\u2022 W20: good reconstruction (~25 dB)', 8, False, GREEN, PP_ALIGN.LEFT)
label(s, explain_x, fg_win_y + Inches(0.6), Inches(4.0), Inches(0.2),
      '\u2022 Full dataset: poor (~10 dB PSNR)', 8, True, RED, PP_ALIGN.LEFT)


# ─── Save ───
out = '/Users/mahesh/Desktop/master-thesis/presentations/arch_experiment.pptx'
prs.save(out)
print('Saved:', out, '| Slides:', len(prs.slides))
