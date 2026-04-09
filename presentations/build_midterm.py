"""
Mid-Term Presentation Builder
Neural Network-Based Compression of Spatio-Temporal Scientific Data
BTU Cottbus-Senftenberg
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
EMU_W = 12192000
EMU_H = 6858000

# Colors
BTU_DARK = RGBColor(0x4A, 0x4A, 0x4A)
BTU_BLUE = RGBColor(0x00, 0x6E, 0xB8)
BTU_GREEN = RGBColor(0x5C, 0xB8, 0x5C)
BTU_ORANGE = RGBColor(0xE8, 0x8D, 0x2A)
BTU_RED = RGBColor(0xDC, 0x35, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RELU_COLOR = RGBColor(0xFF, 0x69, 0x34)  # Orange-red for ReLU
INPUT_COLOR = RGBColor(0x00, 0x6E, 0xB8)  # Blue for input
HIDDEN_COLOR = RGBColor(0x5C, 0xB8, 0x5C)  # Green for hidden
OUTPUT_COLOR = RGBColor(0xE8, 0x8D, 0x2A)  # Orange for output

# Paths
BASE = '/Users/mahesh/Desktop/master-thesis'
PRES = os.path.join(BASE, 'presentations')
RESULTS = os.path.join(BASE, 'results')
LOGO = os.path.join(PRES, 'btu_logo.png')
DESIGN = os.path.join(PRES, 'btu_design.jpg')

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H


# ─── Helpers ───
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                          bold_flags=None, font_name='Calibri', line_spacing=1.2):
    """lines: list of strings, bold_flags: list of bools"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1) + 2)
        if bold_flags and i < len(bold_flags):
            p.font.bold = bold_flags[i]
    return txBox


def add_logo(slide, left=None, top=Inches(0.2), height=Inches(0.55)):
    if left is None:
        left = EMU_W - Inches(3.0)
    slide.shapes.add_picture(LOGO, left, top, height=height)


# ─── Architecture diagram helpers (from arch_experiment) ───
BLUE_F = RGBColor(0x00, 0x6E, 0xB8)
BLUE_T = RGBColor(0x4D, 0xA6, 0xE8)
BLUE_S = RGBColor(0x00, 0x4C, 0x8C)
ORANGE_F = RGBColor(0xE8, 0x8D, 0x2A)
ORANGE_T = RGBColor(0xFF, 0xB3, 0x47)
ORANGE_S = RGBColor(0xCC, 0x77, 0x22)
LAYER_GREEN = RGBColor(0x2E, 0x8B, 0x57)
LAYER_PURPLE = RGBColor(0x6A, 0x0D, 0xAD)
RELU_COL = RGBColor(0xFF, 0x69, 0x34)
RED_COL = RGBColor(0xDC, 0x35, 0x45)
GREEN_COL = RGBColor(0x2E, 0x8B, 0x57)
LOSS_BG_COL = RGBColor(0xFF, 0xF8, 0xE1)
LOSS_BD_COL = RGBColor(0xE8, 0x8D, 0x2A)
CALLOUT_BG_COL = RGBColor(0xFD, 0xEB, 0xD0)
WARN_BG_COL = RGBColor(0xFD, 0xEC, 0xEA)
ARCH_GRAY = RGBColor(0x88, 0x88, 0x88)


def abox(slide, l, t, w, h, fill, text='', sz=10, tc=WHITE, bold=True, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if text:
        s.text_frame.word_wrap = True
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = s.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = tc
        p.font.name = 'Calibri'
    return s


def arbox(slide, l, t, w, h, fill, text='', sz=10, tc=WHITE, bold=True, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    if text:
        s.text_frame.word_wrap = True
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = s.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = tc
        p.font.name = 'Calibri'
    return s


def alabel(slide, l, t, w, h, text, sz=10, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tb


def aarrow(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, Inches(0.3))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
    s.line.fill.background()
    return s


def ahline(slide, l, t, w):
    abox(slide, l, t, w, Pt(1.5), DARK_GRAY)


def avline(slide, l, t, h):
    abox(slide, l, t, Pt(1.5), h, DARK_GRAY)


def abars(slide, cx, cy, n, h, color=LAYER_GREEN, bw=Inches(0.1), gap=Inches(0.04)):
    total = n * bw + (n - 1) * gap
    xs = cx - total // 2
    ys = cy - h // 2
    for i in range(n):
        bx = xs + i * (bw + gap)
        abox(slide, bx, ys, bw, h, color)
    return xs, xs + total


def abracket_top(slide, x1, x2, y, text, sz=9):
    bh = Inches(0.12)
    mid = (x1 + x2) // 2
    avline(slide, x1, y - bh, bh)
    avline(slide, x2, y - bh, bh)
    ahline(slide, x1, y - bh, x2 - x1)
    avline(slide, mid, y - bh - Inches(0.06), Inches(0.06))
    alabel(slide, x1 - Inches(0.5), y - bh - Inches(0.35), x2 - x1 + Inches(1.0), Inches(0.25),
           text, sz=sz, bold=True)


def abracket_bot(slide, x1, x2, y, text, sz=8):
    bh = Inches(0.1)
    avline(slide, x1, y, bh)
    avline(slide, x2, y, bh)
    ahline(slide, x1, y + bh, x2 - x1)
    alabel(slide, x1 - Inches(0.2), y + bh + Inches(0.02), x2 - x1 + Inches(0.4), Inches(0.25),
           text, sz=sz, bold=True)


def atable(slide, left, top, w, h, data, hdr_color, hl_cols=None, fsz=8):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, w, h)
    tbl = ts.table
    cw = w // cols
    for c in range(cols):
        tbl.columns[c].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(fsz)
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = (r == 0)
                p.font.color.rgb = WHITE if r == 0 else (RED_COL if (hl_cols and c in hl_cols) else DARK_GRAY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hdr_color if r == 0 else (WHITE if r % 2 else LIGHT_GRAY)


def add_footer(slide, text="Mid-Term Presentation | Mahesh Sadupalli | BTU Cottbus"):
    add_textbox(slide, Inches(0), EMU_H - Inches(0.45), EMU_W, Inches(0.4),
                text, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def add_slide_title(slide, title, subtitle=None):
    # Blue accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BTU_BLUE
    bar.line.fill.background()

    add_textbox(slide, Inches(0.5), Inches(0.3), EMU_W - Inches(1), Inches(0.7),
                title, font_size=28, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.95), EMU_W - Inches(1), Inches(0.4),
                    subtitle, font_size=14, color=MED_GRAY, alignment=PP_ALIGN.LEFT)
    add_logo(slide)
    add_footer(slide)


def add_rect_box(slide, left, top, width, height, text, fill_color, text_color=WHITE,
                 font_size=12, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    # Vertical center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_rounded_box(slide, left, top, width, height, text, fill_color, text_color=WHITE,
                    font_size=12, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    return shape


def add_arrow(slide, left, top, width, height=Inches(0)):
    """Add a right-pointing arrow connector"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BTU_BLUE
    shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, height, rows, cols, data, header_color=BTU_BLUE,
              col_widths=None, font_size=11):
    """data: list of lists. First row is header."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Calibri'
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.color.rgb = WHITE if r == 0 else DARK_GRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


# ═══════════════════════════════════════════════
# SLIDE 1: COVER PAGE
# ═══════════════════════════════════════════════
slide = add_blank_slide()

# White background (default)
# BTU design image - left side
slide.shapes.add_picture(DESIGN, Inches(0), Inches(0), width=Inches(6.5), height=EMU_H)

# BTU logo - top right
slide.shapes.add_picture(LOGO, EMU_W - Inches(3.5), Inches(0.4), height=Inches(0.7))

# Title text - right side, centered vertically
add_textbox(slide, Inches(6.0), Inches(1.8), Inches(6.8), Inches(1.2),
            "Neural Network-Based Compression of\nSpatio-Temporal Scientific Data",
            font_size=24, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(6.0), Inches(3.2), Inches(6.8), Inches(0.5),
            "Mid-Term Presentation", font_size=18, color=BTU_BLUE, alignment=PP_ALIGN.CENTER)

# Author info
add_multiline_textbox(slide, Inches(6.0), Inches(4.0), Inches(6.8), Inches(2.0),
                      ["Mahesh Sadupalli",
                       "Mtr. Nr: 4034007  |  M.Sc. Artificial Intelligence",
                       "",
                       "Supervisor: Prof. Dr.-Ing. Michael Oevermann",
                       "Mentor: M.Sc. Abhishek Dhiman",
                       "",
                       "BTU Cottbus-Senftenberg",
                       "April 2026"],
                      font_size=14, color=DARK_GRAY,
                      bold_flags=[True, False, False, False, False, False, False])


# ═══════════════════════════════════════════════
# SLIDE 2: OUTLINE
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Outline")

outline_items = [
    ("01", "Motivation & Problem Statement"),
    ("02", "Dataset Overview"),
    ("03", "Methodology: Implicit Neural Representations"),
    ("04", "Offline (Batch) Training — Results"),
    ("05", "Online (Streaming) Training — Results"),
    ("06", "Comparison: Offline vs Online"),
    ("07", "Conclusion & Future Work"),
]

for i, (num, item) in enumerate(outline_items):
    y = Inches(1.6) + Inches(0.68) * i
    # Accent bar on left
    accent = BTU_BLUE if i in [3, 4] else RGBColor(0xCC, 0xCC, 0xCC)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.8), y + Inches(0.08), Inches(0.12), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Item text - LEFT aligned
    txt_color = DARK_GRAY if i in [3, 4] else MED_GRAY
    add_textbox(slide, Inches(2.2), y, Inches(9.0), Inches(0.48),
                num + "    " + item, font_size=17, bold=(i in [3, 4]), color=txt_color,
                alignment=PP_ALIGN.LEFT)

    # Subtle line separator
    if i < len(outline_items) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(1.8), y + Inches(0.58), Inches(9.4), Inches(0.008))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        line.line.fill.background()


# ═══════════════════════════════════════════════
# SLIDE 3: MOTIVATION & PROBLEM
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Motivation & Problem Statement")

# Left column: The problem
add_rounded_box(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5),
                "The Problem", BTU_BLUE, WHITE, 16, True)

problem_lines = [
    "Scientific simulations generate massive data",
    "  Vortex shedding CFD: 794 MB for 300 timesteps",
    "  Climate models: terabytes per simulation run",
    "Storage and transfer become bottlenecks",
    "Traditional compression: limited for scientific data",
]
add_multiline_textbox(slide, Inches(0.7), Inches(2.2), Inches(5.4), Inches(2.5),
                      problem_lines, font_size=14, color=DARK_GRAY,
                      alignment=PP_ALIGN.LEFT, bold_flags=[True, False, False, True, True])

# Right column: Our approach
add_rounded_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.5),
                "Our Approach", BTU_GREEN, WHITE, 16, True)

approach_lines = [
    "Implicit Neural Representations (INRs)",
    "  Neural network learns: f(x,y,z,t) → field variables",
    "  Model weights = compressed representation",
    "Extreme compression ratios: up to ~4,700:1",
    "Domain-agnostic framework",
]
add_multiline_textbox(slide, Inches(7.2), Inches(2.2), Inches(5.4), Inches(2.5),
                      approach_lines, font_size=14, color=DARK_GRAY,
                      alignment=PP_ALIGN.LEFT, bold_flags=[True, False, False, True, True])

# Bottom: key question
add_rounded_box(slide, Inches(2.0), Inches(5.5), Inches(9.3), Inches(0.65),
                "Key Question: Can we compress streaming data online without losing earlier information?",
                RGBColor(0xE8, 0x8D, 0x2A), WHITE, 15, True)


# ═══════════════════════════════════════════════
# SLIDE 4: DATASET
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Dataset Overview", "Vortex Shedding CFD Simulation (Validation Case Study)")

# Dataset stats table
data = [
    ["Property", "Value"],
    ["Total Samples", "7,919,100"],
    ["Timesteps", "300"],
    ["Spatial Points / Timestep", "~26,397"],
    ["Input Features", "x, y, z, t  (4D coordinates)"],
    ["Output Features", "Vx, Vy, Pressure, TKE  (4 field variables)"],
    ["File Size (CSV)", "794 MB"],
    ["Normalization", "Min-max per column → [0, 1]"],
]
add_table(slide, Inches(1.2), Inches(1.5), Inches(6.0), Inches(3.5),
          len(data), 2, data, font_size=12,
          col_widths=[Inches(2.5), Inches(3.5)])

# Right side: data flow diagram
add_rounded_box(slide, Inches(8.0), Inches(1.6), Inches(4.0), Inches(0.6),
                "CSV (No Header, 8 Columns)", RGBColor(0x88, 0x88, 0x88), WHITE, 12)

add_arrow(slide, Inches(9.5), Inches(2.35), Inches(1.0))

add_rounded_box(slide, Inches(8.0), Inches(2.8), Inches(4.0), Inches(0.6),
                "PyArrow Fast Loading", BTU_BLUE, WHITE, 12)

add_arrow(slide, Inches(9.5), Inches(3.55), Inches(1.0))

add_rounded_box(slide, Inches(8.0), Inches(4.0), Inches(1.8), Inches(0.6),
                "Inputs\n(x,y,z,t)", INPUT_COLOR, WHITE, 11)
add_rounded_box(slide, Inches(10.2), Inches(4.0), Inches(1.8), Inches(0.6),
                "Targets\n(Vx,Vy,P,TKE)", OUTPUT_COLOR, WHITE, 11)

add_arrow(slide, Inches(9.5), Inches(4.75), Inches(1.0))

add_rounded_box(slide, Inches(8.0), Inches(5.2), Inches(4.0), Inches(0.6),
                "Min-Max Normalization → [0, 1]", BTU_GREEN, WHITE, 12)

# Flow field image (if available)
viz_path = os.path.join(RESULTS, 'batch_learning/base_model_offline/base_flow_visualization.png')
if os.path.exists(viz_path):
    pass  # Skip large viz, keep slide clean


# ═══════════════════════════════════════════════
# SLIDE 5: OFFLINE INR ARCHITECTURE (PPT native)
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Offline (Batch) Model Architecture")
abox(slide, 0, Inches(0.9), EMU_W, EMU_H - Inches(0.9), LIGHT_GRAY)

cy = Inches(3.3)
cx, ctop, cw, ch = Inches(0.5), Inches(1.9), Inches(1.8), Inches(2.2)
d = Inches(0.18)

# Input cube
tp = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, cx, ctop - d, cw, d + Pt(2))
tp.fill.solid(); tp.fill.fore_color.rgb = BLUE_T; tp.line.color.rgb = BLUE_S; tp.line.width = Pt(0.8)
abox(slide, cx + cw, ctop, d, ch, BLUE_S, border=BLUE_S)
abox(slide, cx, ctop, cw, ch, BLUE_F, 'Input\nData', 13, WHITE, border=BLUE_S)

atable(slide, Inches(0.3), Inches(4.4), Inches(2.2), Inches(0.7),
       [['X', 'Y', 'Z', 'Time'], ['1.01', '2.01', '5.02', '0.01'], ['N...', 'N...', 'N...', 'N...']],
       BLUE_F, fsz=7)

aarrow(slide, Inches(2.65), cy - Inches(0.15), Inches(0.55))

l1l, l1r = abars(slide, Inches(3.75), cy, 3, Inches(0.9), LAYER_PURPLE)
l2l, l2r = abars(slide, Inches(4.85), cy, 6, Inches(2.4), LAYER_GREEN)
l3l, l3r = abars(slide, Inches(6.15), cy, 6, Inches(2.4), LAYER_GREEN)
l4l, l4r = abars(slide, Inches(7.45), cy, 5, Inches(1.7), LAYER_GREEN)
l5l, l5r = abars(slide, Inches(8.55), cy, 3, Inches(0.9), LAYER_PURPLE)

for x in [l1r, l2r, l3r, l4r]:
    aarrow(slide, x + Inches(0.03), cy - Inches(0.1), Inches(0.28))

for xc, lbl in [(Inches(3.75), '4'), (Inches(4.85), '64'), (Inches(6.15), '64'),
                 (Inches(7.45), '32'), (Inches(8.55), '4')]:
    alabel(slide, xc - Inches(0.2), cy + Inches(1.35), Inches(0.4), Inches(0.25), lbl, sz=11, bold=True)

abracket_top(slide, l2l, l4r, cy - Inches(1.3), 'Hidden Layers: ReLU Activation', sz=10)
abracket_bot(slide, l1l, l1r, cy + Inches(0.5), 'Input')
abracket_bot(slide, l2l, l3r, cy + Inches(1.25), 'Hidden Layers (ReLU)')
abracket_bot(slide, l4l, l4r, cy + Inches(0.9), 'Hidden')
abracket_bot(slide, l5l, l5r, cy + Inches(0.5), 'Output')

for rx in [Inches(4.25), Inches(5.45), Inches(6.75)]:
    alabel(slide, rx, cy - Inches(1.35), Inches(0.5), Inches(0.2), 'ReLU', sz=7, bold=True, color=RELU_COL)
alabel(slide, Inches(7.95), cy - Inches(1.3), Inches(0.8), Inches(0.18), 'No Activation', sz=6, color=ARCH_GRAY)

aarrow(slide, l5r + Inches(0.1), cy - Inches(0.15), Inches(0.55))

ox = Inches(9.7)
tp2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, ox, ctop - d, cw, d + Pt(2))
tp2.fill.solid(); tp2.fill.fore_color.rgb = ORANGE_T; tp2.line.color.rgb = ORANGE_S; tp2.line.width = Pt(0.8)
abox(slide, ox + cw, ctop, d, ch, ORANGE_S, border=ORANGE_S)
abox(slide, ox, ctop, cw, ch, ORANGE_F, 'Predicted\nData', 13, WHITE, border=ORANGE_S)

atable(slide, Inches(9.5), Inches(4.4), Inches(2.2), Inches(0.7),
       [['\u0056\u0302x', '\u0056\u0302y', '\u0050\u0302', 'TKE\u0302'],
        ['1.05', '-0.32', '0.95', '12.3'], ['N...', 'N...', 'N...', 'N...']],
       ORANGE_F, hl_cols=[0, 1, 2, 3], fsz=7)

arbox(slide, Inches(4.0), Inches(5.5), Inches(4.0), Inches(0.55),
      LOSS_BG_COL, 'Loss: MSE  |  Metrics: PSNR, SSIM', 9, DARK_GRAY, True, LOSS_BD_COL)

arbox(slide, Inches(1.5), Inches(6.3), Inches(9.0), Inches(0.45),
      CALLOUT_BG_COL,
      '4 \u2192 64 \u2192 64 \u2192 32 \u2192 4  |  6,692 params  |  26.1 KB  |  CR ~4,734:1',
      9, RGBColor(0x8B, 0x40, 0x00), True, RGBColor(0xE8, 0x8D, 0x2A))


# ═══════════════════════════════════════════════
# SLIDE 7: COMPRESSION CONCEPT
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Compression via INR", "Store model weights instead of full dataset")

compress_img = os.path.join(PRES, 'diagrams', 'compression.png')
if os.path.exists(compress_img):
    slide.shapes.add_picture(compress_img, Inches(0.5), Inches(1.4), width=Inches(12.0))

# Compression table at bottom
cr_data = [
    ["Model", "Parameters", "Model Size", "Dataset Size", "Compression Ratio"],
    ["BaseCompressor", "6,692", "26.1 KB", "120.8 MB", "~4,734 : 1"],
    ["MediumCompressor", "14,644", "57.2 KB", "120.8 MB", "~2,165 : 1"],
    ["LargeCompressor", "25,668", "100.3 KB", "120.8 MB", "~1,233 : 1"],
]
add_table(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.8),
          len(cr_data), 5, cr_data, font_size=12,
          col_widths=[Inches(2.2), Inches(1.8), Inches(1.8), Inches(2.0), Inches(2.5)])


# ═══════════════════════════════════════════════
# SLIDE 8: OFFLINE TRAINING APPROACH
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Offline (Batch) Training", "Train on entire dataset for 150 epochs")

# Training flow diagram using PPT shapes
add_rounded_box(slide, Inches(0.5), Inches(2.0), Inches(2.2), Inches(1.0),
                "Full Dataset\n7.9M samples", INPUT_COLOR, WHITE, 13, True)
add_arrow(slide, Inches(2.8), Inches(2.35), Inches(0.8))

add_rounded_box(slide, Inches(3.7), Inches(2.0), Inches(2.0), Inches(1.0),
                "DataLoader\nbatch_size=512\nshuffle=True", BTU_DARK, WHITE, 11, True)
add_arrow(slide, Inches(5.8), Inches(2.35), Inches(0.8))

add_rounded_box(slide, Inches(6.7), Inches(2.0), Inches(2.2), Inches(1.0),
                "MLP Model\nfθ(x,y,z,t)", HIDDEN_COLOR, WHITE, 13, True)
add_arrow(slide, Inches(9.0), Inches(2.35), Inches(0.8))

add_rounded_box(slide, Inches(9.9), Inches(2.0), Inches(2.2), Inches(1.0),
                "Predictions\n(V̂x, V̂y, P̂, k̂)", OUTPUT_COLOR, WHITE, 13, True)

# Feedback loop
add_textbox(slide, Inches(4.5), Inches(3.2), Inches(4.5), Inches(0.4),
            "← MSE Loss → Adam Optimizer → Update θ → Repeat 150 epochs",
            font_size=12, bold=True, color=BTU_BLUE)

# Hyperparameters table
hyper_data = [
    ["Parameter", "Value"],
    ["Epochs", "150"],
    ["Batch Size", "512"],
    ["Optimizer", "Adam (lr=0.001)"],
    ["Loss Function", "MSE Loss"],
    ["Activation", "ReLU (hidden layers only)"],
]
add_table(slide, Inches(3.5), Inches(4.0), Inches(6.0), Inches(2.4),
          len(hyper_data), 2, hyper_data, font_size=12,
          col_widths=[Inches(2.5), Inches(3.5)])


# ═══════════════════════════════════════════════
# SLIDE 8: OFFLINE RESULTS
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Offline Results", "Full-dataset evaluation after 150 epochs")

# Results table
results_data = [
    ["Model", "PSNR (dB) ↑", "SSIM ↑", "Rel. Error (%) ↓", "Compression Ratio"],
    ["Base (6.7K)", "32.15", "0.9550", "4.41", "~4,734:1"],
    ["Medium (14.6K)", "33.58", "0.9583", "3.74", "~2,165:1"],
    ["Large (25.7K)", "35.99", "0.9856", "2.83", "~1,233:1"],
]
add_table(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.0),
          len(results_data), 5, results_data, font_size=14,
          col_widths=[Inches(2.0), Inches(2.0), Inches(1.8), Inches(2.3), Inches(2.2)])

# Add comparison plots (sized to fit within slide)
comp_path = os.path.join(RESULTS, 'batch_learning/comparison_offline/offline_training_comparison.png')
if os.path.exists(comp_path):
    slide.shapes.add_picture(comp_path, Inches(0.5), Inches(3.8),
                             width=Inches(5.8), height=Inches(3.0))

eval_path = os.path.join(RESULTS, 'batch_learning/comparison_offline/offline_evaluation_comparison.png')
if os.path.exists(eval_path):
    slide.shapes.add_picture(eval_path, Inches(6.8), Inches(3.8),
                             width=Inches(5.8), height=Inches(3.0))


# ═══════════════════════════════════════════════
# SLIDE 9: OFFLINE FLOW VISUALIZATION
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Offline: Flow Field Reconstruction", "Original vs Predicted vs Absolute Error (Large Model)")

viz_path = os.path.join(RESULTS, 'batch_learning/large_model_offline/large_flow_visualization.png')
if os.path.exists(viz_path):
    slide.shapes.add_picture(viz_path, Inches(0.5), Inches(1.4),
                             width=Inches(11.0))


# ═══════════════════════════════════════════════
# SLIDE 11: ONLINE ARCHITECTURE (PPT native)
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Online (Streaming) Model Architecture")
abox(slide, 0, Inches(0.9), EMU_W, EMU_H - Inches(0.9), LIGHT_GRAY)

cy = Inches(2.9)
cx, ctop, cw, ch = Inches(0.4), Inches(1.6), Inches(1.7), Inches(2.0)
d = Inches(0.15)

# Input cube
tp = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, cx, ctop - d, cw, d + Pt(2))
tp.fill.solid(); tp.fill.fore_color.rgb = BLUE_T; tp.line.color.rgb = BLUE_S; tp.line.width = Pt(0.8)
abox(slide, cx + cw, ctop, d, ch, BLUE_S, border=BLUE_S)
abox(slide, cx, ctop, cw, ch, BLUE_F, 'Streaming\nData', 12, WHITE, border=BLUE_S)

atable(slide, Inches(0.3), Inches(3.85), Inches(2.0), Inches(0.6),
       [['X', 'Y', 'Z', 'Time'], ['1.01', '2.01', '5.02', '0.01'], ['N...', 'N...', 'N...', 'N...']],
       BLUE_F, fsz=6)

aarrow(slide, Inches(2.3), cy - Inches(0.12), Inches(0.4))

# Temporal windowing area
win_area_x = Inches(2.85)
win_area_w = Inches(2.6)
win_area_y = Inches(1.45)
win_area_h = Inches(2.8)
arbox(slide, win_area_x, win_area_y, win_area_w, win_area_h,
      RGBColor(0xE3, 0xF2, 0xFD), '', border=RGBColor(0x90, 0xCA, 0xF9))

alabel(slide, win_area_x, win_area_y + Inches(0.05), win_area_w, Inches(0.25),
       'Temporal Windowing', 9, True)

win_w = Inches(0.42)
win_h = Inches(0.85)
win_row_x = win_area_x + Inches(0.12)
win_row_y = win_area_y + Inches(0.4)

win_data = [
    ('W1', RGBColor(0x90, 0xCA, 0xF9), DARK_GRAY),
    ('W2', RGBColor(0x42, 0xA5, 0xF5), WHITE),
    ('...', RGBColor(0xF5, 0xF5, 0xF5), ARCH_GRAY),
    ('W19', RGBColor(0x1E, 0x88, 0xE5), WHITE),
    ('W20', RGBColor(0x15, 0x65, 0xC0), WHITE),
]
actual_gap = (win_area_w - Inches(0.24) - 5 * win_w) // 4
for i, (wlbl, wfill, wtc) in enumerate(win_data):
    wx = win_row_x + i * (win_w + actual_gap)
    abox(slide, wx, win_row_y, win_w, win_h, wfill, wlbl, 7 if len(wlbl) > 2 else 8, wtc, True,
         RGBColor(0x0D, 0x47, 0xA1) if wlbl != '...' else RGBColor(0xCC, 0xCC, 0xCC))

alabel(slide, win_area_x, win_row_y + win_h + Inches(0.08), win_area_w, Inches(0.2),
       '15 timesteps per window', 7, False, ARCH_GRAY)
alabel(slide, win_area_x, win_row_y + win_h + Inches(0.28), win_area_w, Inches(0.2),
       '20 windows \u00d7 100 epochs = 2,000 total', 7, True, RGBColor(0x15, 0x65, 0xC0))

aarrow(slide, win_area_x + win_area_w + Inches(0.05), cy - Inches(0.12), Inches(0.35))

l2l, l2r = abars(slide, Inches(6.5), cy, 6, Inches(2.2), LAYER_GREEN)
l3l, l3r = abars(slide, Inches(7.5), cy, 6, Inches(2.2), LAYER_GREEN)
l4l, l4r = abars(slide, Inches(8.5), cy, 5, Inches(1.5), LAYER_GREEN)
l5l, l5r = abars(slide, Inches(9.35), cy, 3, Inches(0.8), LAYER_PURPLE)

for x in [l2r, l3r, l4r]:
    aarrow(slide, x + Inches(0.03), cy - Inches(0.1), Inches(0.2))

for xc, lbl in [(Inches(6.5), '64'), (Inches(7.5), '64'), (Inches(8.5), '32'), (Inches(9.35), '4')]:
    alabel(slide, xc - Inches(0.2), cy + Inches(1.2), Inches(0.4), Inches(0.25), lbl, sz=10, bold=True)

abracket_top(slide, l2l, l4r, cy - Inches(1.2), 'Hidden Layers: ReLU', sz=9)
abracket_bot(slide, win_area_x, win_area_x + win_area_w, cy + Inches(0.6), 'Input (4)')
abracket_bot(slide, l2l, l3r, cy + Inches(1.15), 'Hidden (ReLU)')
abracket_bot(slide, l4l, l4r, cy + Inches(0.8), 'Hidden')
abracket_bot(slide, l5l, l5r, cy + Inches(0.45), 'Output')

for rx in [Inches(6.9), Inches(7.95)]:
    alabel(slide, rx, cy - Inches(1.2), Inches(0.45), Inches(0.18), 'ReLU', sz=7, bold=True, color=RELU_COL)
alabel(slide, Inches(8.85), cy - Inches(1.15), Inches(0.7), Inches(0.16), 'No Activation', sz=6, color=ARCH_GRAY)

aarrow(slide, l5r + Inches(0.08), cy - Inches(0.12), Inches(0.4))

ox = Inches(10.2)
tp2 = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, ox, ctop - d, cw, d + Pt(2))
tp2.fill.solid(); tp2.fill.fore_color.rgb = ORANGE_T; tp2.line.color.rgb = ORANGE_S; tp2.line.width = Pt(0.8)
abox(slide, ox + cw, ctop, d, ch, ORANGE_S, border=ORANGE_S)
abox(slide, ox, ctop, cw, ch, ORANGE_F, 'Predicted\nData', 12, WHITE, border=ORANGE_S)

atable(slide, Inches(10.1), Inches(3.85), Inches(2.0), Inches(0.6),
       [['\u0056\u0302x', '\u0056\u0302y', '\u0050\u0302', 'TKE\u0302'],
        ['1.05', '-0.32', '0.95', '12.3'], ['N...', 'N...', 'N...', 'N...']],
       ORANGE_F, hl_cols=[0, 1, 2, 3], fsz=6)

arbox(slide, Inches(4.8), Inches(4.65), Inches(3.5), Inches(0.45),
      LOSS_BG_COL, 'Loss: MSE  |  Metrics: PSNR, SSIM', 8, DARK_GRAY, True, LOSS_BD_COL)

# Catastrophic forgetting element
forget_y = Inches(5.35)
arbox(slide, Inches(0.4), forget_y, Inches(12.0), Inches(1.6),
      WARN_BG_COL, '', border=RED_COL)
alabel(slide, Inches(0.6), forget_y + Inches(0.05), Inches(11.5), Inches(0.25),
       '\u26A0  Catastrophic Forgetting', 11, True, RED_COL)

fg_win_y = forget_y + Inches(0.35)
fg_win_w = Inches(0.42)
fg_win_h = Inches(0.5)
fg_gap = Inches(0.08)
fg_start_x = Inches(0.7)

fg_windows = [
    ('W1', RGBColor(0xEF, 0x9A, 0x9A), RED_COL, True),
    ('W2', RGBColor(0xEF, 0x9A, 0x9A), RED_COL, True),
    ('W3', RGBColor(0xEF, 0x9A, 0x9A), RED_COL, True),
    ('W4', RGBColor(0xEF, 0x9A, 0x9A), RED_COL, True),
    ('W5', RGBColor(0xEF, 0x9A, 0x9A), RED_COL, True),
    ('...', RGBColor(0xF5, 0xF5, 0xF5), ARCH_GRAY, False),
    ('W18', RGBColor(0xEF, 0x9A, 0x9A), RED_COL, True),
    ('W19', RGBColor(0xEF, 0x9A, 0x9A), RED_COL, True),
    ('W20', RGBColor(0xA5, 0xD6, 0xA7), GREEN_COL, False),
]
for i, (wlbl, wfill, wtc, forgotten) in enumerate(fg_windows):
    wx = fg_start_x + i * (fg_win_w + fg_gap)
    abox(slide, wx, fg_win_y, fg_win_w, fg_win_h, wfill, wlbl, 7, wtc, True,
         RED_COL if forgotten else GREEN_COL)
    if wlbl != '...':
        sym = '\u2717' if forgotten else '\u2713'
        alabel(slide, wx, fg_win_y + fg_win_h + Pt(2), fg_win_w, Inches(0.15),
               sym, 9, True, RED_COL if forgotten else GREEN_COL)

explain_x = fg_start_x + 9 * (fg_win_w + fg_gap) + Inches(0.2)
alabel(slide, explain_x, fg_win_y - Inches(0.05), Inches(4.0), Inches(0.2),
       'After training on W20:', 9, True, DARK_GRAY, PP_ALIGN.LEFT)
alabel(slide, explain_x, fg_win_y + Inches(0.2), Inches(4.0), Inches(0.18),
       '\u2022 W1\u2013W19: knowledge overwritten', 8, False, RED_COL, PP_ALIGN.LEFT)
alabel(slide, explain_x, fg_win_y + Inches(0.4), Inches(4.0), Inches(0.18),
       '\u2022 W20: good reconstruction (~25 dB)', 8, False, GREEN_COL, PP_ALIGN.LEFT)
alabel(slide, explain_x, fg_win_y + Inches(0.6), Inches(4.0), Inches(0.18),
       '\u2022 Full dataset: poor (~10 dB PSNR)', 8, True, RED_COL, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════
# SLIDE 11: ONLINE RESULTS
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Online Results", "Per-window streaming metrics vs Full-dataset evaluation")

# Last window results
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
            "Last Window Performance (Window 20)", font_size=14, bold=True, color=BTU_GREEN)

last_window = [
    ["Model", "PSNR (dB)", "SSIM", "Rel. Error (%)", "Time (s)"],
    ["Base", "23.99", "0.8740", "11.93", "43.56"],
    ["Medium", "24.40", "0.8806", "11.37", "66.62"],
    ["Large", "27.45", "0.9017", "8.01", "81.52"],
]
add_table(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(1.8),
          len(last_window), 5, last_window, BTU_GREEN, font_size=12,
          col_widths=[Inches(1.1), Inches(1.1), Inches(1.0), Inches(1.3), Inches(1.3)])

# Full dataset results
add_textbox(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
            "Full-Dataset Evaluation (after all 20 windows)", font_size=14, bold=True, color=BTU_RED)

full_eval = [
    ["Model", "PSNR (dB)", "SSIM", "Rel. Error (%)"],
    ["Base", "11.97", "0.7551", "44.92"],
    ["Medium", "12.70", "0.7599", "41.30"],
    ["Large", "9.67", "0.6679", "58.57"],
]
add_table(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(1.8),
          len(full_eval), 4, full_eval, BTU_RED, font_size=12,
          col_widths=[Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.6)])

# Comparison plots
comp_online = os.path.join(RESULTS, 'continual_learning/comparison_online/online_training_comparison.png')
if os.path.exists(comp_online):
    slide.shapes.add_picture(comp_online, Inches(0.5), Inches(4.0),
                             width=Inches(5.8), height=Inches(3.0))

eval_online = os.path.join(RESULTS, 'continual_learning/comparison_online/online_evaluation_comparison.png')
if os.path.exists(eval_online):
    slide.shapes.add_picture(eval_online, Inches(6.8), Inches(4.0),
                             width=Inches(5.8), height=Inches(3.0))


# ═══════════════════════════════════════════════
# SLIDE 12: ONLINE FLOW VISUALIZATION
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Online: Flow Field Reconstruction", "Full-dataset evaluation — showing catastrophic forgetting (Large Model)")

viz_online = os.path.join(RESULTS, 'continual_learning/large_model_online/large_online_visualization.png')
if os.path.exists(viz_online):
    slide.shapes.add_picture(viz_online, Inches(0.5), Inches(1.4), width=Inches(11.0))


# ═══════════════════════════════════════════════
# SLIDE 14: COMPARISON OFFLINE vs ONLINE (rich diagram + table)
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Comparison: Offline vs Online", "Same models, different training strategies")

# Embed the rich comparison diagram
comp_img = os.path.join(PRES, 'diagrams', 'comparison.png')
if os.path.exists(comp_img):
    slide.shapes.add_picture(comp_img, Inches(0.3), Inches(1.3), width=Inches(12.0))

# Compact comparison table below
comp_data = [
    ["", "Offline (Batch)", "Online (Streaming)"],
    ["Base PSNR", "32.15 dB", "11.97 dB"],
    ["Medium PSNR", "33.58 dB", "12.70 dB"],
    ["Large PSNR", "35.99 dB", "9.67 dB"],
]
add_table(slide, Inches(3.0), Inches(5.0), Inches(7.3), Inches(1.8),
          len(comp_data), 3, comp_data, font_size=13,
          col_widths=[Inches(2.3), Inches(2.5), Inches(2.5)])

# Key insight
add_rounded_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.55),
                "Gap: ~20 dB PSNR loss in online mode — the central research challenge",
                BTU_ORANGE, WHITE, 13, True)


# ═══════════════════════════════════════════════
# SLIDE: COMPRESSION RATIO CALCULATION
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Compression Ratio Calculation")

cr_calc_img = os.path.join(PRES, 'diagrams', 'compression_calc.png')
if os.path.exists(cr_calc_img):
    slide.shapes.add_picture(cr_calc_img, Inches(0.5), Inches(1.4), width=Inches(12.0))


# ═══════════════════════════════════════════════
# SLIDE: CONCLUSION
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Conclusion")

conclusions = [
    ("INRs achieve extreme compression", "~4,700:1 with BaseCompressor (26 KB for 241 MB data)"),
    ("Offline training: high quality", "PSNR up to 35.99 dB, SSIM 0.9856 with Large model"),
    ("Online training: streaming capable", "But suffers severe catastrophic forgetting (~20 dB drop)"),
    ("Larger models ≠ always better", "Large model forgets MORE in online mode (9.67 vs 11.97 dB)"),
    ("Framework is domain-agnostic", "Applicable to climate, molecular dynamics, structural mechanics, etc."),
]

for i, (title, detail) in enumerate(conclusions):
    y = Inches(1.5) + Inches(1.05) * i
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BTU_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(1.7), y - Inches(0.05), Inches(10.0), Inches(0.35),
                title, font_size=16, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1.7), y + Inches(0.32), Inches(10.0), Inches(0.35),
                detail, font_size=13, color=MED_GRAY, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════
# SLIDE 15: FUTURE WORK
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "Future Work")

# Future work items with status
future_items = [
    ("Continual Learning Strategies", "Experience Replay, EWC, LwF to mitigate catastrophic forgetting", BTU_BLUE),
    ("Experience Replay Variants", "ER-Scaled (weighted replay) and ER-Aggressive (full replay)", BTU_BLUE),
    ("Alternative Architectures", "Autoencoders (Linear, Conv2D) for comparison", BTU_GREEN),
    ("Hyperparameter Optimization", "Window sizes, buffer sizes, learning rates", MED_GRAY),
    ("Benchmark on Other Domains", "Climate data, molecular dynamics, structural mechanics", MED_GRAY),
]

for i, (title, detail, color) in enumerate(future_items):
    y = Inches(1.5) + Inches(1.05) * i

    # Status bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y, Inches(0.15), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_textbox(slide, Inches(1.4), y - Inches(0.05), Inches(10.5), Inches(0.35),
                title, font_size=16, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1.4), y + Inches(0.32), Inches(10.5), Inches(0.35),
                detail, font_size=13, color=MED_GRAY, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════
# SLIDE 16: THANK YOU
# ═══════════════════════════════════════════════
slide = add_blank_slide()

# BTU design on left
slide.shapes.add_picture(DESIGN, Inches(0), Inches(0), width=Inches(6.5), height=EMU_H)
slide.shapes.add_picture(LOGO, EMU_W - Inches(3.5), Inches(0.4), height=Inches(0.7))

add_textbox(slide, Inches(6.0), Inches(2.5), Inches(6.8), Inches(1.0),
            "Thank You!", font_size=36, bold=True, color=DARK_GRAY)

add_textbox(slide, Inches(6.0), Inches(3.5), Inches(6.8), Inches(0.5),
            "Questions & Discussion", font_size=20, color=BTU_BLUE)

add_multiline_textbox(slide, Inches(6.0), Inches(4.5), Inches(6.8), Inches(1.5),
                      ["Mahesh Sadupalli",
                       "sadupmah@b-tu.de",
                       "",
                       "Supervisor: Prof. Dr.-Ing. Michael Oevermann",
                       "Mentor: M.Sc. Abhishek Dhiman"],
                      font_size=13, color=DARK_GRAY,
                      bold_flags=[True, False, False, False, False])


# ─── Save ───
output_path = os.path.join(PRES, 'midterm_presentation.pptx')
prs.save(output_path)
print('Saved to:', output_path)
print('Total slides:', len(prs.slides))
