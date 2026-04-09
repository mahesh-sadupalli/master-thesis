#!/usr/bin/env python3
"""Update the Offline_vs_Online_Comparison.pptx with medium model results and comparison images."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy
from lxml import etree
import os

PPTX_PATH = "/Users/mahesh/Desktop/master-thesis/presentations/Offline_vs_Online_Comparison.pptx"
RESULTS_DIR = "/Users/mahesh/Desktop/master-thesis/results"

prs = Presentation(PPTX_PATH)
slide_width = prs.slide_width   # 12192000 EMU (standard widescreen)
slide_height = prs.slide_height # 6858000 EMU


def copy_cell_format(src_cell, dst_cell):
    """Copy formatting from source cell to destination cell."""
    # Copy text formatting
    if src_cell.text_frame.paragraphs:
        src_para = src_cell.text_frame.paragraphs[0]
        dst_para = dst_cell.text_frame.paragraphs[0]
        if src_para.runs:
            src_run = src_para.runs[0]
            for dst_run in dst_para.runs:
                dst_run.font.size = src_run.font.size
                dst_run.font.bold = src_run.font.bold
                if src_run.font.color and src_run.font.color.rgb:
                    dst_run.font.color.rgb = src_run.font.color.rgb
        dst_para.alignment = src_para.alignment


def set_cell_text(cell, text, font_size=Pt(14), bold=False, alignment=PP_ALIGN.CENTER, color=None):
    """Set cell text with formatting."""
    cell.text = ""
    para = cell.text_frame.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_column_to_table(table, col_idx, header_text, row_data, ref_col=1):
    """Add a column to an existing table by duplicating XML and inserting data."""
    # python-pptx doesn't support adding columns directly, so we modify XML
    tbl = table._tbl
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    # Get gridCol elements
    tblGrid = tbl.find('.//a:tblGrid', nsmap)
    gridCols = tblGrid.findall('a:gridCol', nsmap)

    # Calculate new column width (take from reference column)
    ref_gridCol = gridCols[ref_col]
    new_width = int(ref_gridCol.get('w'))

    # Create new gridCol
    new_gridCol = etree.SubElement(tblGrid, '{http://schemas.openxmlformats.org/drawingml/2006/main}gridCol')
    new_gridCol.set('w', str(new_width))

    # Insert at correct position
    if col_idx < len(gridCols):
        tblGrid.remove(new_gridCol)
        tblGrid.insert(col_idx, new_gridCol)

    # Add cells to each row
    rows = tbl.findall('.//a:tr', nsmap)
    for row_idx, tr in enumerate(rows):
        tcs = tr.findall('a:tc', nsmap)
        # Clone the reference cell
        ref_tc = tcs[ref_col]
        new_tc = copy.deepcopy(ref_tc)

        # Clear text in new cell
        for p in new_tc.findall('.//a:p', nsmap):
            for r in p.findall('a:r', nsmap):
                t = r.find('a:t', nsmap)
                if t is not None:
                    if row_idx == 0:
                        t.text = header_text
                    elif row_idx - 1 < len(row_data):
                        t.text = str(row_data[row_idx - 1])

        # Insert at correct position
        if col_idx < len(tcs):
            tr.insert(col_idx, new_tc)
        else:
            tr.append(new_tc)


def add_row_to_table(table, row_data, ref_row=1):
    """Add a row to an existing table."""
    tbl = table._tbl
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    rows = tbl.findall('.//a:tr', nsmap)
    ref_tr = rows[ref_row]
    new_tr = copy.deepcopy(ref_tr)

    tcs = new_tr.findall('a:tc', nsmap)
    for i, tc in enumerate(tcs):
        for p in tc.findall('.//a:p', nsmap):
            for r in p.findall('a:r', nsmap):
                t = r.find('a:t', nsmap)
                if t is not None and i < len(row_data):
                    t.text = str(row_data[i])

    # Insert after last data row (before any totals)
    ref_tr.addnext(new_tr)


def add_text_box(slide, left, top, width, height, text, font_size=Pt(14), bold=False, alignment=PP_ALIGN.CENTER, color=None):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_image_slide(prs, title_text, image_path, subtitle_text=None):
    """Add a new slide with a title and full-width image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(11.2), Inches(0.6),
                 title_text, font_size=Pt(28), bold=True, alignment=PP_ALIGN.LEFT)

    # Image
    if os.path.exists(image_path):
        img_top = Inches(0.85)
        img_height = Inches(5.8) if not subtitle_text else Inches(5.4)
        pic = slide.shapes.add_picture(image_path, Inches(0.5), img_top, Inches(11.2), img_height)

    # Subtitle
    if subtitle_text:
        add_text_box(slide, Inches(0.5), Inches(6.4), Inches(11.2), Inches(0.4),
                     subtitle_text, font_size=Pt(12), bold=False, alignment=PP_ALIGN.CENTER,
                     color=RGBColor(0x66, 0x66, 0x66))

    return slide


# ============================================================
# SLIDE 2: Compression Ratio — Offline Training
# Add Medium CR column
# ============================================================
print("Updating Slide 2: Compression Ratio — Offline Training")
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_table:
        table = shape.table
        # Add Medium column between Base and Large
        # Current: Method | Original Size | Base CR | Large CR
        # Target:  Method | Original Size | Base CR | Medium CR | Large CR
        add_column_to_table(table, 3, "Medium CR", ["4,146 : 1", "13,630 : 1", "60 KB"], ref_col=2)
    if shape.has_text_frame and "Model sizes" in shape.text:
        shape.text_frame.paragraphs[0].runs[0].text = "Model sizes: Base .pth = 30 KB  |  Medium .pth = 60 KB  |  Large .pth = 104 KB"

# ============================================================
# SLIDE 3: Offline Training Results
# Add Medium row and update metrics
# ============================================================
print("Updating Slide 3: Offline Training Results")
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_table:
        table = shape.table
        # Update Base row (row 1) with new metrics
        tbl = table._tbl
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        rows = tbl.findall('.//a:tr', nsmap)

        # Update Base row values
        base_data = ['Base', '4→64→64→32→4', '6,692', '32.15', '0.9550', '4.41']
        tcs = rows[1].findall('a:tc', nsmap)
        for i, tc in enumerate(tcs):
            for p in tc.findall('.//a:p', nsmap):
                for r in p.findall('a:r', nsmap):
                    t = r.find('a:t', nsmap)
                    if t is not None and i < len(base_data):
                        t.text = base_data[i]

        # Update Large row values
        large_data = ['Large', '4→128→128→64→4', '25,668', '35.99', '0.9856', '2.83']
        tcs = rows[2].findall('a:tc', nsmap)
        for i, tc in enumerate(tcs):
            for p in tc.findall('.//a:p', nsmap):
                for r in p.findall('a:r', nsmap):
                    t = r.find('a:t', nsmap)
                    if t is not None and i < len(large_data):
                        t.text = large_data[i]

        # Add Medium row between Base and Large
        add_row_to_table(table, ['Medium', '4→96→96→48→4', '14,644', '33.58', '0.9583', '3.74'], ref_row=1)

    # Update callout text
    if shape.has_text_frame and "Large model achieves" in shape.text:
        shape.text_frame.paragraphs[0].runs[0].text = "Large model achieves PSNR 35.99 dB and SSIM 0.9856 — excellent reconstruction quality"

# ============================================================
# SLIDE 7: Compression Ratio — Online Training
# Add Medium CR column
# ============================================================
print("Updating Slide 7: Compression Ratio — Online Training")
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_table:
        table = shape.table
        add_column_to_table(table, 3, "Medium CR", ["4,146 : 1", "13,630 : 1", "60 KB"], ref_col=2)
    if shape.has_text_frame and "Model sizes" in shape.text:
        shape.text_frame.paragraphs[0].runs[0].text = "Model sizes: Base .pth = 30 KB  |  Medium .pth = 60 KB  |  Large .pth = 104 KB"

# ============================================================
# SLIDE 8: Online Streaming Results
# Add Medium row and update metrics
# ============================================================
print("Updating Slide 8: Online Streaming Results")
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.has_table:
        table = shape.table
        tbl = table._tbl
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        rows = tbl.findall('.//a:tr', nsmap)

        # Update Base row
        base_data = ['Base', '23.99', '0.8740', '11.97', '0.7551']
        tcs = rows[1].findall('a:tc', nsmap)
        for i, tc in enumerate(tcs):
            for p in tc.findall('.//a:p', nsmap):
                for r in p.findall('a:r', nsmap):
                    t = r.find('a:t', nsmap)
                    if t is not None and i < len(base_data):
                        t.text = base_data[i]

        # Update Large row
        large_data = ['Large', '27.45', '0.9017', '9.67', '0.6679']
        tcs = rows[2].findall('a:tc', nsmap)
        for i, tc in enumerate(tcs):
            for p in tc.findall('.//a:p', nsmap):
                for r in p.findall('a:r', nsmap):
                    t = r.find('a:t', nsmap)
                    if t is not None and i < len(large_data):
                        t.text = large_data[i]

        # Add Medium row
        add_row_to_table(table, ['Medium', '24.40', '0.8806', '12.70', '0.7599'], ref_row=1)

# ============================================================
# SLIDE 12: Offline vs Online — Head-to-Head
# Add Medium columns and update metrics
# ============================================================
print("Updating Slide 12: Offline vs Online — Head-to-Head")
slide12 = prs.slides[11]
for shape in slide12.shapes:
    if shape.has_table:
        table = shape.table
        tbl = table._tbl
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        rows = tbl.findall('.//a:tr', nsmap)

        # Update header: add Medium Offline and Medium Online columns
        # Current: Metric | Base Offline | Base Online | Large Offline | Large Online
        # Target:  Metric | Base Offline | Base Online | Med Offline | Med Online | Large Offline | Large Online

        # Update existing values first
        # Row 1 - PSNR
        psnr_data = ['PSNR (dB)', '32.15', '11.97', '35.99', '9.67']
        tcs = rows[1].findall('a:tc', nsmap)
        for i, tc in enumerate(tcs):
            for p in tc.findall('.//a:p', nsmap):
                for r in p.findall('a:r', nsmap):
                    t = r.find('a:t', nsmap)
                    if t is not None and i < len(psnr_data):
                        t.text = psnr_data[i]

        # Row 2 - SSIM
        ssim_data = ['SSIM', '0.9550', '0.7551', '0.9856', '0.6679']
        tcs = rows[2].findall('a:tc', nsmap)
        for i, tc in enumerate(tcs):
            for p in tc.findall('.//a:p', nsmap):
                for r in p.findall('a:r', nsmap):
                    t = r.find('a:t', nsmap)
                    if t is not None and i < len(ssim_data):
                        t.text = ssim_data[i]

        # Add Medium columns (between Base Online and Large Offline)
        add_column_to_table(table, 3, "Med Offline", ["33.58", "0.9583"], ref_col=1)
        add_column_to_table(table, 4, "Med Online", ["12.70", "0.7599"], ref_col=2)

    # Update callout
    if shape.has_text_frame and "PSNR Drop" in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "PSNR Drop" in run.text:
                    run.text = "PSNR Drop:  Base 32.15→11.97 (−20.18 dB)  |  Med 33.58→12.70 (−20.88 dB)  |  Large 35.99→9.67 (−26.32 dB)"
                elif "SSIM Drop" in run.text:
                    run.text = "SSIM Drop:  Base 0.955→0.755 (−0.200)  |  Med 0.958→0.760 (−0.198)  |  Large 0.986→0.668 (−0.318)"

# ============================================================
# SLIDE 13: Key Takeaways - update text
# ============================================================
print("Updating Slide 13: Key Takeaways")
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text
            if "8,277:1" in text or "2,379:1" in text:
                for run in para.runs:
                    run.text = "2.  Compression ratios: Base 27,395:1  |  Medium 13,241:1  |  Large 7,713:1 — identical for both modes"
            elif "PSNR >30" in text:
                for run in para.runs:
                    if "PSNR" in run.text:
                        run.text = "1.  Offline training achieves excellent quality: Base 32.15 dB, Medium 33.58 dB, Large 35.99 dB PSNR"

# ============================================================
# ADD NEW SLIDES: Offline comparison images (4 metrics)
# ============================================================
print("Adding offline comparison image slides")

# Insert after slide 4 (offline training progress)
# We'll add them at the end and note the intended position
offline_comparison_images = [
    ("Offline Comparison — Training Progress", f"{RESULTS_DIR}/Offline_comparison/offline_training_comparison.png",
     "Training convergence comparison across Base, Medium, and Large models"),
    ("Offline Comparison — Evaluation Metrics", f"{RESULTS_DIR}/Offline_comparison/offline_evaluation_comparison.png",
     "Final evaluation metrics comparison across all three model sizes"),
    ("Offline Comparison — Loss", f"{RESULTS_DIR}/Offline_comparison/loss_comparison.png",
     "Loss convergence comparison across model sizes"),
    ("Offline Comparison — PSNR", f"{RESULTS_DIR}/Offline_comparison/psnr_comparison.png",
     "PSNR progression comparison across model sizes"),
    ("Offline Comparison — SSIM", f"{RESULTS_DIR}/Offline_comparison/ssim_comparison.png",
     "SSIM convergence comparison across model sizes"),
    ("Offline Comparison — Relative Error", f"{RESULTS_DIR}/Offline_comparison/relative_error_comparison.png",
     "Relative error comparison across model sizes"),
]

for title, img_path, subtitle in offline_comparison_images:
    add_image_slide(prs, title, img_path, subtitle)

# ============================================================
# ADD NEW SLIDES: Online comparison images (4 metrics)
# ============================================================
print("Adding online comparison image slides")

online_comparison_images = [
    ("Online Comparison — Training Progress", f"{RESULTS_DIR}/comparison_online/online_training_comparison.png",
     "Online training convergence comparison across Base, Medium, and Large models"),
    ("Online Comparison — Evaluation Metrics", f"{RESULTS_DIR}/comparison_online/online_evaluation_comparison.png",
     "Online evaluation metrics comparison across all three model sizes"),
    ("Online Comparison — Loss", f"{RESULTS_DIR}/comparison_online/loss_comparison.png",
     "Online loss comparison across model sizes"),
    ("Online Comparison — PSNR", f"{RESULTS_DIR}/comparison_online/psnr_comparison.png",
     "Online PSNR comparison across model sizes"),
    ("Online Comparison — SSIM", f"{RESULTS_DIR}/comparison_online/ssim_comparison.png",
     "Online SSIM comparison across model sizes"),
    ("Online Comparison — Relative Error", f"{RESULTS_DIR}/comparison_online/relative_error_comparison.png",
     "Online relative error comparison across model sizes"),
]

for title, img_path, subtitle in online_comparison_images:
    add_image_slide(prs, title, img_path, subtitle)

# ============================================================
# ADD: Medium model visualization slides
# ============================================================
print("Adding medium model visualization slides")

medium_vis_images = [
    ("Medium Model — Training Progress (Offline)", f"{RESULTS_DIR}/medium_model_offline/medium_model_training_progress.png", None),
    ("Medium Model — Flow Field Reconstruction (Offline)", f"{RESULTS_DIR}/medium_model_offline/medium_offline_visualization.png", None),
    ("Medium Model — Training Progress (Online)", f"{RESULTS_DIR}/medium_model_online/medium_model_online_progress.png", None),
    ("Medium Model — Flow Field Reconstruction (Online)", f"{RESULTS_DIR}/medium_model_online/medium_online_visualization.png", None),
]

for title, img_path, subtitle in medium_vis_images:
    add_image_slide(prs, title, img_path, subtitle)

# ============================================================
# SAVE
# ============================================================
output_path = PPTX_PATH
prs.save(output_path)
print(f"\nSaved updated presentation to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
